import os
import timeit

from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt
import pandas as pd
from pymongo import MongoClient
import base64
import io
import json

def PowerPoint_IV(wafer_id):
    """
    This function takes a path to a txt file in argument and creates a Powerpoint where are stored the plots between voltage (V) and current (A)
    :param <str> original_file: Path to your .txt file where measurements are stored
    :return: None
    """
    start_time = timeit.default_timer()
    if os.path.exists(f"PowerPointFiles\{wafer_id} plots_IV.pptx"):
        return
    client = MongoClient('mongodb://localhost:27017/')
    db = client['Measurements']
    collection = db["Wafers"]
    wafer = collection.find_one({"wafer_id": wafer_id})

    # creating directory if it doesn't exist
    if not os.path.exists("plots"):
        os.makedirs("plots")
    if not os.path.exists("PowerPointFiles"):
        os.makedirs("PowerPointFiles")

    prs = Presentation()
    colors = [
        '#FF0000',  # Rouge
        '#00FF00',  # Vert
        '#0000FF',  # Bleu
        '#FFFF00',  # Jaune
        '#FF00FF',  # Magenta
        '#00FFFF',  # Cyan
        '#800000',  # Marron
        '#808000',  # Olive
        '#008000',  # Vert foncé
        '#008080',  # Vert d'eau
        '#000080',  # Bleu marine
        '#800080',  # Violet
        '#7F7F7F',  # Gris
        '#FF6600',  # Orange
        '#663399'  # Rebecca Purple
    ]

    # Creating lists of sheets that already exists in the excel files
    list_of_sheets = []
    df_dict = {}  # using dictionary to store dataframes by coord keys

    # Processing files
    for structure in wafer["structures"]:
        for matrix in structure["matrices"]:
            coord = "("+matrix["coordinates"]["x"]+','+matrix["coordinates"]["y"]+')'

            if coord not in df_dict.keys():
                df_dict[coord] = pd.DataFrame()

            testdeviceID = structure["structure_id"]
            voltages = ['Voltage (V)']
            I = ['I (A)']
            for double in matrix["results"]["I"]["Values"]:
                voltages.append(double["V"])
                I.append(double["I"])

            new_df = pd.DataFrame(list(zip(voltages, I)), columns=[testdeviceID, testdeviceID+" "])
            df_dict[coord] = df_dict[coord].merge(new_df, left_index=True, right_index=True, how = 'outer')

    for coord, df in df_dict.items():
        cols = [col.rstrip() for col in df.columns if col.endswith(' ')]
        plt.figure()
        color_index = 0
        for col in cols:

            label_x = df[col].iloc[0]
            label_y = "Current Value (A)"

            df[col] = df[col].iloc[1:]
            df[col + ' '] = df[col + ' '].iloc[1:]
            df[col + ' '] = abs(df[col + ' '])

            plt.plot(df[col], df[col+" "], str(colors[color_index%15]), label = col)

            plt.xlabel(label_x)
            plt.ylabel(label_y)
            plt.title(wafer_id + " " + coord)
            plt.legend()
            plt.grid(True)

            color_index += 1

        plt.savefig("plots\\" + wafer_id + coord + '.png')
        plt.close()

        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        left = Inches(1)
        top = Inches(1)

        slide.shapes.add_picture("plots\\" + wafer_id + coord + '.png', left, top)

        prs.save(f"PowerPointFiles\{wafer_id} plots_IV.pptx")
    end_time = timeit.default_timer()
    print(f"I-V PowerPoint successfully created for {wafer_id} in {end_time-start_time} seconds!")

def PowerPoint_JV(wafer_id):
    """
        This function takes a path to a txt file in argument and creates a Powerpoint where are stored the plots between voltage (V) and current (A)
        :param <str> original_file: Path to your .txt file where measurements are stored
        :return: None
        """
    start_time = timeit.default_timer()
    if os.path.exists(f"PowerPointFiles\{wafer_id} plots_JV.pptx"):
        return
    client = MongoClient('mongodb://localhost:27017/')
    db = client['Measurements']
    collection = db["Wafers"]
    wafer = collection.find_one({"wafer_id": wafer_id})

    # creating directory if it doesn't exist
    if not os.path.exists("plots"):
        os.makedirs("plots")
    if not os.path.exists("PowerPointFiles"):
        os.makedirs("PowerPointFiles")

    prs = Presentation()
    colors = [
        '#FF0000',  # Rouge
        '#00FF00',  # Vert
        '#0000FF',  # Bleu
        '#FFFF00',  # Jaune
        '#FF00FF',  # Magenta
        '#00FFFF',  # Cyan
        '#800000',  # Marron
        '#808000',  # Olive
        '#008000',  # Vert foncé
        '#008080',  # Vert d'eau
        '#000080',  # Bleu marine
        '#800080',  # Violet
        '#7F7F7F',  # Gris
        '#FF6600',  # Orange
        '#663399'  # Rebecca Purple
    ]

    # Creating lists of sheets that already exists in the excel files
    list_of_sheets = []
    df_dict = {}  # using dictionary to store dataframes by coord keys

    # Processing files
    for structure in wafer["structures"]:
        for matrix in structure["matrices"]:
            coord = "(" + matrix["coordinates"]["x"] + ',' + matrix["coordinates"]["y"] + ')'

            if coord not in df_dict.keys():
                df_dict[coord] = pd.DataFrame()

            testdeviceID = structure["structure_id"]
            voltages = ['Voltage (V)']
            J = ['J (A/cm^2)']
            for double in matrix["results"]["J"]["Values"]:
                voltages.append(double["V"])
                J.append(double["J"])

            new_df = pd.DataFrame(list(zip(voltages, J)), columns=[testdeviceID, testdeviceID + " "])
            df_dict[coord] = df_dict[coord].merge(new_df, left_index=True, right_index=True, how='outer')

    for coord, df in df_dict.items():
        cols = [col.rstrip() for col in df.columns if col.endswith(' ')]
        plt.figure()
        color_index = 0
        for col in cols:
            label_x = df[col].iloc[0]
            label_y = "Current density (A/cm^2)"

            df[col] = df[col].iloc[1:]
            df[col + ' '] = df[col + ' '].iloc[1:]
            df[col + ' '] = abs(df[col + ' '])

            plt.plot(df[col], df[col + " "], str(colors[color_index % 15]), label=col)

            plt.xlabel(label_x)
            plt.ylabel(label_y)
            plt.title(wafer_id + " " + coord)
            plt.legend()
            plt.grid(True)

            color_index += 1

        plt.savefig("plots\\" + wafer_id + coord + '.png')
        plt.close()

        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        left = Inches(1)
        top = Inches(1)

        slide.shapes.add_picture("plots\\" + wafer_id + coord + '.png', left, top)

        prs.save(f"PowerPointFiles\{wafer_id} plots_JV.pptx")
    end_time = timeit.default_timer()
    print(f"J-V PowerPoint successfully created for {wafer_id} in {end_time-start_time} seconds!")

def writeppt(wafer_id):
    """
    This function creates a PowerPoint presentation where plots of different data types are stored.
    :param <str> wafer_id: The id of the wafer
    :return: None
    """
    print("Starting Powerpoint")
    data_types = ['I', 'J', 'It']
    y_values_dict = {'I': ['I'], 'J': ['J'], 'It': ['It']}

    client = MongoClient('mongodb://localhost:27017/')
    db = client['Measurements']
    collection = db["Wafers"]
    wafer = collection.find_one({"wafer_id": wafer_id})

    # creating directory if it doesn't exist
    if not os.path.exists("plots"):
        os.makedirs("plots")
    if not os.path.exists("PowerPointFiles"):
        os.makedirs("PowerPointFiles")

    colors = [
        '#FF0000',  # Rouge
        '#00FF00',  # Vert
        '#0000FF',  # Bleu
        '#FFFF00',  # Jaune
        '#FF00FF',  # Magenta
        '#00FFFF',  # Cyan
        '#800000',  # Marron
        '#808000',  # Olive
        '#008000',  # Vert foncé
        '#008080',  # Vert d'eau
        '#000080',  # Bleu marine
        '#800080',  # Violet
        '#7F7F7F',  # Gris
        '#FF6600',  # Orange
        '#663399'  # Rebecca Purple
    ]

    for data_type in data_types:
        y_values = y_values_dict[data_type]
        data_label = f'{data_type} Values'

        start_time = timeit.default_timer()
        if os.path.exists(f"PowerPointFiles\{wafer_id} plots_{data_type}.pptx"):
            continue

        df_dict = {}  # using dictionary to store dataframes by coord keys

        # Processing files
        for structure in wafer["structures"]:
            for matrix in structure["matrices"]:
                #print(f"Matrix results: {matrix['results']}")
                # check if the data_type is in results, if not, skip this loop iteration
                if data_type not in matrix["results"]:
                    continue

                coord = "(" + matrix["coordinates"]["x"] + ',' + matrix["coordinates"]["y"] + ')'

                if coord not in df_dict.keys():
                    df_dict[coord] = pd.DataFrame()

                testdeviceID = structure["structure_id"]
                data_columns = [f'{data_type} ({unit})' for unit in y_values]
                data_values = [[] for _ in range(len(y_values) + 1)]
                data_values[0] = ['Voltage (V)']

                for double in matrix["results"][data_type]["Values"]:
                    data_values[0].append(double["V"])
                    #print(f"Y values:{y_values}")
                    for idx, unit in enumerate(y_values, 1):
                        data_values[idx].append(double[unit])

                columns = [testdeviceID] + [testdeviceID + " " + unit for unit in y_values]
                new_df = pd.DataFrame(list(zip(*data_values)), columns=columns)
                df_dict[coord] = df_dict[coord].merge(new_df, left_index=True, right_index=True, how='outer')

        prs = Presentation()

        for coord, df in df_dict.items():
            for unit in y_values:
                cols = [col.rstrip() for col in df.columns if not col.endswith(' ' + unit)]
                if not cols:  # If no columns for this unit, skip
                    continue

                plt.figure()
                color_index = 0
                for col in cols:
                    label_x = df[col].iloc[0]
                    label_y = f"{data_type} Value ({unit})"

                    df[col] = df[col].iloc[1:]
                    df[col + ' ' + unit] = df[col + ' ' + unit].iloc[1:]
                    df[col + ' ' + unit] = abs(df[col + ' ' + unit])

                    plt.plot(df[col], df[col + ' ' + unit], str(colors[color_index % 15]), label=col)

                    plt.xlabel(label_x)
                    plt.ylabel(label_y)
                    plt.title(f"{wafer_id} {coord}")
                    plt.legend()
                    plt.grid(True)

                    color_index += 1

                plt.savefig(f"plots\\{wafer_id}{coord}_{data_type}_{unit}.png")
                plt.close()

                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)

                left = Inches(1)
                top = Inches(1)

                slide.shapes.add_picture(f"plots\\{wafer_id}{coord}_{data_type}_{unit}.png", left, top)

        if len(prs.slides) > 0:
            prs.save(f"PowerPointFiles\\{wafer_id}_plots_{data_type}.pptx")
            end_time = timeit.default_timer()

def ppt_structure(wafer_id=str, structure_ids=list, file_name=str):
    data_types = ['I', 'J', 'It']
    y_values_dict = {'I': ['I'], 'J': ['J'], 'It': ['It']}

    client = MongoClient('mongodb://localhost:27017/')
    db = client['Measurements']
    collection = db["Wafers"]
    wafer = collection.find_one({"wafer_id": wafer_id})

    # creating directory if it doesn't exist
    if not os.path.exists("plots"):
        os.makedirs("plots")
    if not os.path.exists("PowerPointFiles"):
        os.makedirs("PowerPointFiles")

    colors = [
        '#FF0000',  # Rouge
        '#00FF00',  # Vert
        '#0000FF',  # Bleu
        '#FFFF00',  # Jaune
        '#FF00FF',  # Magenta
        '#00FFFF',  # Cyan
        '#800000',  # Marron
        '#808000',  # Olive
        '#008000',  # Vert foncé
        '#008080',  # Vert d'eau
        '#000080',  # Bleu marine
        '#800080',  # Violet
        '#7F7F7F',  # Gris
        '#FF6600',  # Orange
        '#663399'  # Rebecca Purple
    ]

    for data_type in data_types:
        y_values = y_values_dict[data_type]
        data_label = f'{data_type} Values'

        if os.path.exists(f"PowerPointFiles\\{file_name}.pptx"):
            continue

        df_dict = {}  # using dictionary to store dataframes by coord keys

        # Processing files
        for structure in wafer["structures"]:
            if structure["structure_id"] in structure_ids:
                for matrix in structure["matrices"]:
                    # print(f"Matrix results: {matrix['results']}")
                    # check if the data_type is in results, if not, skip this loop iteration
                    if data_type not in matrix["results"]:
                        continue

                    coord = "(" + matrix["coordinates"]["x"] + ',' + matrix["coordinates"]["y"] + ')'

                    if coord not in df_dict.keys():
                        df_dict[coord] = pd.DataFrame()

                    testdeviceID = structure["structure_id"]
                    data_columns = [f'{data_type} ({unit})' for unit in y_values]
                    data_values = [[] for _ in range(len(y_values) + 1)]
                    data_values[0] = ['Voltage (V)']

                    for double in matrix["results"][data_type]["Values"]:
                        data_values[0].append(double["V"])
                        # print(f"Y values:{y_values}")
                        for idx, unit in enumerate(y_values, 1):
                            data_values[idx].append(double[unit])

                    columns = [testdeviceID] + [testdeviceID + " " + unit for unit in y_values]
                    new_df = pd.DataFrame(list(zip(*data_values)), columns=columns)
                    df_dict[coord] = df_dict[coord].merge(new_df, left_index=True, right_index=True, how='outer')

        prs = Presentation()

        for coord, df in df_dict.items():
            for unit in y_values:
                cols = [col.rstrip() for col in df.columns if not col.endswith(' ' + unit)]
                if not cols:  # If no columns for this unit, skip
                    continue

                plt.figure()
                color_index = 0
                for col in cols:
                    label_x = df[col].iloc[0]
                    label_y = f"{data_type} Value ({unit})"

                    df[col] = df[col].iloc[1:]
                    df[col + ' ' + unit] = df[col + ' ' + unit].iloc[1:]
                    df[col + ' ' + unit] = abs(df[col + ' ' + unit])

                    plt.plot(df[col], df[col + ' ' + unit], str(colors[color_index % 15]), label=col)

                    plt.xlabel(label_x)
                    plt.ylabel(label_y)
                    plt.title(f"{wafer_id} {coord}")
                    plt.legend()
                    plt.grid(True)

                    color_index += 1

                plt.savefig(f"plots\\{file_name}_{wafer_id}{coord}_{data_type}_{unit}.png")
                plt.close()

                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)

                left = Inches(1)
                top = Inches(1)

                slide.shapes.add_picture(f"plots\\{file_name}_{wafer_id}{coord}_{data_type}_{unit}.png", left, top)

        if len(prs.slides) > 0:
            prs.save(f"PowerPointFiles\\{file_name}.pptx")
    return

def ppt_matrix(wafer_id=str, coordinates=str):
    x = coordinates.split(",")[0][1:]
    y = coordinates.split(",")[-1][:-1]

    data_types = ['I', 'J', 'It']
    y_values_dict = {'I': ['I'], 'J': ['J'], 'It': ['It']}

    client = MongoClient('mongodb://localhost:27017/')
    db = client['Measurements']
    collection = db["Wafers"]
    wafer = collection.find_one({"wafer_id": wafer_id})

    figs = []

    colors = [
        '#FF0000',  # Rouge
        '#00FF00',  # Vert
        '#0000FF',  # Bleu
        '#FFFF00',  # Jaune
        '#FF00FF',  # Magenta
        '#00FFFF',  # Cyan
        '#800000',  # Marron
        '#808000',  # Olive
        '#008000',  # Vert foncé
        '#008080',  # Vert d'eau
        '#000080',  # Bleu marine
        '#800080',  # Violet
        '#7F7F7F',  # Gris
        '#FF6600',  # Orange
        '#663399'  # Rebecca Purple
    ]

    for data_type in data_types:
        y_values = y_values_dict[data_type]
        data_label = f'{data_type} Values'

        df_dict = {}

        for structure in wafer["structures"]:
            for matrix in structure["matrices"]:
                if matrix["coordinates"]["x"] == x and matrix["coordinates"]["y"] == y:
                    if coordinates not in df_dict.keys():
                        df_dict[coordinates] = pd.DataFrame()

                    testdeviceID = structure["structure_id"]
                    data_columns = [f'{data_type} ({unit})' for unit in y_values]
                    data_values = [[] for _ in range(len(y_values) + 1)]
                    data_values[0] = ['Voltage (V)']

                    if data_type in matrix["results"]:
                        for double in matrix["results"][data_type]["Values"]:
                            data_values[0].append(double["V"])
                            for idx, unit in enumerate(y_values, 1):
                                data_values[idx].append(double[unit])

                        columns = [testdeviceID] + [testdeviceID + " " + unit for unit in y_values]
                        new_df = pd.DataFrame(list(zip(*data_values)), columns=columns)
                        df_dict[coordinates] = df_dict[coordinates].merge(new_df, left_index=True, right_index=True, how='outer')

        for coord, df in df_dict.items():
            for unit in y_values:
                cols = [col.rstrip() for col in df.columns if not col.endswith(' ' + unit)]
                if not cols:  # If no columns for this unit, skip
                    continue

                fig, ax = plt.subplots(figsize=(10, 5))  # Crée une figure avec une taille spécifique
                color_index = 0

                for col in cols:
                    if not df[col].empty:
                        label_x = df[col].iloc[0]
                        label_y = f"{data_type} Value ({unit})"

                        df[col] = df[col].iloc[1:]
                        df[col + ' ' + unit] = df[col + ' ' + unit].iloc[1:]
                        df[col + ' ' + unit] = abs(df[col + ' ' + unit])

                        ax.plot(df[col], df[col + ' ' + unit], str(colors[color_index % 15]), label=col)

                        ax.set_xlabel(label_x)
                        ax.set_ylabel(label_y)
                        ax.set_title(f"{wafer_id} {coord}")
                        ax.legend(loc='upper left', bbox_to_anchor=(1, 1), fontsize='small')
                        ax.grid(True)

                        color_index += 1

                plt.subplots_adjust(right=0.7)  # Ajuste la zone de dessin pour laisser de l'espace à la légende
                figs.append(fig)
    return [fig_to_base64(fig) for fig in figs]

def fig_to_base64(fig):
    fig_file = io.BytesIO()
    fig.savefig(fig_file, format='png')
    fig_file.seek(0)
    fig_png_base64 = base64.b64encode(fig_file.read())
    return fig_png_base64.decode('utf-8')
